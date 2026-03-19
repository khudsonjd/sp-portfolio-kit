"""
Generate SharePoint Custom Solution Catalog PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x4A, 0x90, 0xD9)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2A, 0x2A, 0x2A)
GREY   = RGBColor(0x55, 0x55, 0x55)
LGREY  = RGBColor(0x99, 0x99, 0x99)

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


# ── Helpers ───────────────────────────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    """Solid-fill the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()   # no border
    return shape


def add_line(slide, x1, y1, x2, y2, color: RGBColor, width_pt=1.5):
    """Add a horizontal rule as a thin rectangle."""
    height = Pt(1.5)
    rect = add_rect(slide, x1, y1, x2 - x1, height, color)
    return rect


def add_textbox(slide, left, top, width, height, text,
                font_name="Calibri", font_size=18, bold=False, italic=False,
                color: RGBColor = DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox, tf


def add_footer(slide):
    """Add 'MDFFS Foundation' label + navy square in bottom-right on every slide."""
    # Small navy square (logo stand-in)
    sq_size = Inches(0.12)
    sq_left = W - Inches(1.75)
    sq_top  = H - Inches(0.28)
    add_rect(slide, sq_left, sq_top, sq_size, sq_size, NAVY)

    # Label
    tb, tf = add_textbox(
        slide,
        left   = sq_left + sq_size + Inches(0.05),
        top    = H - Inches(0.30),
        width  = Inches(1.5),
        height = Inches(0.25),
        text   = "MDFFS Foundation",
        font_size = 9,
        color     = BLUE,
    )


def add_top_bar(slide):
    """6pt-tall horizontal bar across the top, solid navy."""
    add_rect(slide, 0, 0, W, Pt(6), NAVY)


def add_content_header(slide, heading_text):
    """
    Accent vertical bar + heading used on content slides 2-5.
    Returns the y-position just below the heading for body text.
    """
    bar_left  = Inches(0.5)
    bar_top   = Inches(1.1)
    bar_w     = Inches(0.06)
    bar_h     = Inches(0.55)
    add_rect(slide, bar_left, bar_top, bar_w, bar_h, BLUE)

    heading_left = bar_left + bar_w + Inches(0.12)
    tb, tf = add_textbox(
        slide,
        left      = heading_left,
        top       = bar_top - Inches(0.05),
        width     = Inches(11.5),
        height    = Inches(0.65),
        text      = heading_text,
        font_size = 28,
        bold      = True,
        color     = NAVY,
    )
    return bar_top + bar_h + Inches(0.25)


def add_body_paragraph(tf, text, font_size=18, bold=False, color=DARK,
                        space_before_pt=0, add_blank_before=False):
    """Append a paragraph to an existing text frame."""
    if add_blank_before:
        blank = tf.add_paragraph()
        blank.text = ""
        blank.space_before = Pt(space_before_pt)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_hyperlink_run(paragraph, text, url, font_size=16, color=BLUE, underline=True):
    """Add a hyperlinked run to an existing paragraph."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import nsmap

    # Build the rId for the relationship
    slide = paragraph._p.getroottree().getroot()
    # We need the slide part — walk up to find it
    part = paragraph._p

    # Use the run approach via XML directly
    rPr_attrib = {
        "lang": "en-US",
        "dirty": "0",
    }

    # Create <a:r> run element
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_elem = etree.SubElement(paragraph._p, qn("a:r"))

    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "en-US", "dirty": "0"})

    # Solid fill for color
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"),
                               attrib={"val": f"{color.rgb:06X}"})
    # Font size
    rPr.set("sz", str(int(font_size * 100)))
    if underline:
        rPr.set("u", "sng")

    # Hyperlink — add hlinkClick to rPr
    # We need the slide relationship part
    # Walk up the XML tree to find the presentation part
    # Simpler: inject the rId via the slide's part directly
    # We'll use a placeholder and fix it after
    hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"),
                                  attrib={qn("r:id"): "rId_placeholder"})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text

    return r_elem, hlinkClick, url


def fix_hyperlinks(slide, hyperlink_data):
    """
    hyperlink_data: list of (hlinkClick_element, url)
    Add relationships and patch rId values.
    """
    for hlink_elem, url in hyperlink_data:
        rId = slide.part.relate_to(url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True)
        hlink_elem.set(qn("r:id"), rId)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide1(prs):
    """Title slide — dark navy background, centered title."""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_top_bar(slide)
    add_footer(slide)

    cx = W / 2
    rule_w = Inches(4)
    rule_y_top = Inches(2.7)
    rule_y_bot = Inches(3.85)

    # Top decorative rule
    add_line(slide, cx - rule_w // 2, rule_y_top, cx + rule_w // 2, rule_y_top, BLUE)

    # Title text box (two lines)
    title_left = Inches(2)
    title_top  = Inches(2.85)
    title_w    = Inches(9.33)
    title_h    = Inches(1.2)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = txBox.text_frame
    tf.word_wrap = False

    for line_text in ["SharePoint Custom", "Solution Catalog"]:
        p = tf.add_paragraph() if line_text != "SharePoint Custom" else tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line_text
        run.font.name = "Calibri"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Bottom decorative rule
    add_line(slide, cx - rule_w // 2, rule_y_bot, cx + rule_w // 2, rule_y_bot, BLUE)


def build_content_slide(prs, heading, body_lines, font_sizes, bolds=None,
                        colors=None, blank_befores=None, hyperlinks=None):
    """
    Generic content slide (slides 2-5).
    body_lines   : list of strings
    font_sizes   : list of ints (one per line)
    bolds        : list of bools
    colors       : list of RGBColor
    blank_befores: list of bools — insert blank paragraph before this line
    hyperlinks   : list of (line_index, url) or None per line
    """
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT)
    add_top_bar(slide)
    add_footer(slide)

    body_top = add_content_header(slide, heading)

    body_left = Inches(0.68)
    body_w    = Inches(12.0)
    body_h    = H - body_top - Inches(0.4)

    txBox = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    if bolds        is None: bolds        = [False] * len(body_lines)
    if colors       is None: colors       = [DARK]  * len(body_lines)
    if blank_befores is None: blank_befores = [False] * len(body_lines)

    hyperlink_data = []
    first = True
    for i, (text, fsize, bold, color, blank_before) in enumerate(
            zip(body_lines, font_sizes, bolds, colors, blank_befores)):

        if blank_before and not first:
            blank_p = tf.add_paragraph()
            blank_p.text = ""

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(fsize)
        run.font.bold = bold
        run.font.color.rgb = color

        # Check for inline hyperlink replacement (for lines with a URL at the end)
        if hyperlinks and hyperlinks[i]:
            url = hyperlinks[i]
            # Find where the URL starts in the text
            url_start = text.find("http")
            if url_start >= 0:
                pre_text = text[:url_start]
                url_text = text[url_start:]
                # Reset the run text to pre-text only
                run.text = pre_text

                # Add a new run for the URL with hyperlink
                url_run = p.add_run()
                url_run.text = url_text
                url_run.font.name = "Calibri"
                url_run.font.size = Pt(fsize)
                url_run.font.bold = bold
                url_run.font.color.rgb = BLUE
                url_run.font.underline = True

                # Add hyperlink relationship
                rId = slide.part.relate_to(
                    url,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                    is_external=True
                )
                # Inject hlinkClick into the run's rPr
                r_elem = url_run._r
                rPr = r_elem.find(qn("a:rPr"))
                if rPr is None:
                    rPr = etree.SubElement(r_elem, qn("a:rPr"))
                hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"),
                                              attrib={qn("r:id"): rId})

    return slide


def build_slide6(prs):
    """Thank You slide."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_top_bar(slide)

    cx = W / 2
    content_top = Inches(1.9)

    # Heading
    tb, tf = add_textbox(
        slide,
        left      = Inches(2),
        top       = content_top,
        width     = Inches(9.33),
        height    = Inches(0.85),
        text      = "Thank You",
        font_size = 40,
        bold      = True,
        color     = NAVY,
        align     = PP_ALIGN.CENTER,
    )

    # Decorative rule below heading
    rule_y = content_top + Inches(0.9)
    add_line(slide, cx - Inches(2), rule_y, cx + Inches(2), rule_y, BLUE)

    # Name
    add_textbox(
        slide,
        left      = Inches(2),
        top       = rule_y + Inches(0.25),
        width     = Inches(9.33),
        height    = Inches(0.45),
        text      = "Keith Hudson",
        font_size = 20,
        bold      = True,
        color     = NAVY,
        align     = PP_ALIGN.CENTER,
    )

    # Title
    add_textbox(
        slide,
        left      = Inches(2),
        top       = rule_y + Inches(0.75),
        width     = Inches(9.33),
        height    = Inches(0.40),
        text      = "Founder, MDFFS Foundation",
        font_size = 16,
        color     = GREY,
        align     = PP_ALIGN.CENTER,
    )

    # Email with hyperlink
    email_top = rule_y + Inches(1.2)
    email_left = Inches(2)
    email_w    = Inches(9.33)
    email_h    = Inches(0.40)

    txBox = slide.shapes.add_textbox(email_left, email_top, email_w, email_h)
    tf_email = txBox.text_frame
    p = tf_email.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "khudsonjd@gmail.com"
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.color.rgb = BLUE

    # Add hyperlink
    mailto = "mailto:khudsonjd@gmail.com"
    rId = slide.part.relate_to(
        mailto,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True
    )
    r_elem = run._r
    rPr = r_elem.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(r_elem, qn("a:rPr"))
    etree.SubElement(rPr, qn("a:hlinkClick"), attrib={qn("r:id"): rId})

    # Footer — logo square + label
    add_footer(slide)

    # Copyright — bottom left
    add_textbox(
        slide,
        left      = Inches(0.3),
        top       = H - Inches(0.30),
        width     = Inches(5),
        height    = Inches(0.25),
        text      = "© 2026 MDFFS Foundation. All rights reserved.",
        font_size = 9,
        color     = LGREY,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    # Slide 1 — Title
    build_slide1(prs)

    # Slide 2 — Situation
    build_content_slide(
        prs,
        heading    = "Situation",
        body_lines = [
            "SharePoint is a rapid development platform available to everyone in the enterprise."
        ],
        font_sizes  = [18],
    )

    # Slide 3 — Complication
    build_content_slide(
        prs,
        heading    = "Complication",
        body_lines = [
            "SharePoint solutions are often built to meet urgent, unfunded, or \"temporary\" needs — "
            "with plans to sunset them in a few months or rebuild them properly once funding arrives.",
            "Years later, those tools are still running. Your support team is maintaining dozens or "
            "hundreds of custom solutions of uncertain value, with no systematic way to evaluate them.",
        ],
        font_sizes   = [16, 16],
        blank_befores = [False, True],
    )

    # Slide 4 — Resolution
    build_content_slide(
        prs,
        heading    = "Resolution",
        body_lines = [
            "Learn to think like an Enterprise Architect. Apply TOGAF principles to create a "
            "standardized, repeatable process for regularly reviewing your solution portfolio — "
            "informed by relevant organizational and technical context.",
            "Create an automated SharePoint Solution Portfolio Catalog, as explained at "
            "https://github.com/khudsonjd/sp-portfolio-kit",
        ],
        font_sizes   = [16, 16],
        blank_befores = [False, True],
        hyperlinks   = [
            None,
            "https://github.com/khudsonjd/sp-portfolio-kit",
        ],
    )

    # Slide 5 — Resolution (Demo)
    build_content_slide(
        prs,
        heading    = "Resolution",
        body_lines = ["Prototype demo: TAP Catalog"],
        font_sizes  = [32],
        bolds       = [True],
        colors      = [NAVY],
    )

    # Slide 6 — Thank You
    build_slide6(prs)

    out = "SharePoint-Custom-Solution-Catalog.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
